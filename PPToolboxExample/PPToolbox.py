# Powerpoint Toolbox
# Natt Srisutthiyakorn

from pptx import Presentation
from pptx.util import Inches
import os
import numpy as np

# Function
def createPPT(pp_fnTemplate, pp_fnOutput, fig_dir):
#    """
#    Automatically create powerpoint by inserting photos from the folders. This
#    file is currently optimized for the Shell powerpoint template. The function 
#    take the filename or subfolder name a title for a slide.
#    
#    Inputs:
#    - pp_fnTemplate         filename for the powerpoint input 
#                            (generated from SHELL WIZKIT)
#    - pp_fnOutput           filename for the powerpoint output
#    - fig_dir               directory of figures. If there is a subfolder within 
#                            this directory, the function will resize the images and
#                            put them together.
#    
#    Future upgrade:
#    Automatically calculate the size of the figure to make sure that it is within the bound of slide.    
#    
#    Example:
#    pp_dir          = r'\\americas.shell.com\tcs\hou\ua.sepco\proj\mex_exp1\bidround\ndi\sb_ns\misc\powerpoint'
#    pp_fnTemplate   = os.path.join(pp_dir, 'ShellTemplate.pptx') 
#    pp_fnOutput     = os.path.join(pp_dir, '20191102_Geohazard_Jalapeno_Seismic_Test.pptx')
#    fig_dir         = r'\\americas.shell.com\tcs\hou\ua.sepco\proj\mex_exp1\bidround\ndi\sb_ns\misc\snapshots\20190211'
#    
#    createPPT(pp_fnTemplate, pp_fnOutput, fig_dir)   
#    
#    """    
    # Open the powerpoint file
    prs = Presentation(pp_fnTemplate)
    
    # Select Layout #7 (#6 in python) (Title and Content - High Content)
    slide_layout  = prs.slide_layouts[6]

    # List Figure and Folder with in the main folder.
    listFigure = os.listdir(fig_dir)
    listFigure.sort()
    

            
    for file in listFigure:
        
        print(file)
        
        # Shell template specification for wide screen presentation (inch)
        corner_left     = 0.56 # 0.56 
        corner_top      = 1.67 
        total_height    = 5.3
        fig_gap         = 0.3
    
        # Full filename
        fig_fn      = os.path.join(fig_dir, file)  
        
        # Add slide
        slide   = prs.slides.add_slide(slide_layout)
        
        # Check if this is a file or a folder
        if file[-4:] == '.png' or file[-4:] == '.jpg' or file[-4:] == '.tiff' or file[-5:] == '.jpeg':

            # Add the figure
            pic         = slide.shapes.add_picture(fig_fn, Inches(corner_left), Inches(corner_top), height = Inches(total_height))
            
            # Add the title
            title       = slide.shapes.title
            title.text  = file[:-4] 
            
        if os.path.isdir(fig_fn):
            # Number of columns and rows in different case
            listCase    = np.array([1, 4, 6, 9, 12, 16, 20, 25, 30, 36, 42, 49])
            listCol     = np.array([1, 2, 3, 3,  4,  4,  5,  5,  6,  6,  7,  7])
            listRow     = np.array([1, 2, 2, 3,  3,  4,  4,  5,  5,  6,  6,  7])
            listCase    = listCol*listRow
            
            # Find the number of files in the folder
            listFigureSubFolder = os.listdir(fig_fn)
            listFigureSubFolder.sort()
            listFigureSubFolder = [ x for x in listFigureSubFolder if ".db" not in x ]
            
            # Find the number of figure and appropriate number of col and row
            nFig        = len(listFigureSubFolder)
            idxCase     = np.argmax(listCase >= nFig)
            nRow        = listRow[idxCase]
            nCol        = listCol[idxCase]
            
            # Solve for the image location (unit = inch)
            fig_height      = total_height/nRow - fig_gap        
            space_row       = fig_height + fig_gap
            space_col       = 1.74 * fig_height + fig_gap 
            # Assume image width = 1.78 * image height (screen size 2560/1440) 
            
            if nFig > 4:
                corner_left     = 0.15
            
            # Generate location of each figure
            loc_left   = np.arange(0, nCol * space_col, space_col) + corner_left
            loc_top    = np.arange(0, nRow * space_row, space_row) + corner_top
            loc_left, loc_top = np.meshgrid(loc_left, loc_top)
            loc_left   = np.concatenate(loc_left)
            loc_top    = np.concatenate(loc_top)
        
            for count, iFile in enumerate(listFigureSubFolder):
                figSubFolder_fn = os.path.join(fig_fn, iFile)
                pic = slide.shapes.add_picture(figSubFolder_fn, Inches(loc_left[count]), Inches(loc_top[count]), height = Inches(fig_height))

            # Add the title
            title       = slide.shapes.title
            title.text  = file
            
    prs.save(pp_fnOutput)
    
    
#==============================================================================
# Run Example
pp_dir          = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\Python_Toolbox_NS\PPToolboxExample'
pp_fnTemplate   = os.path.join(pp_dir, 'ShellTemplate.pptx') 
pp_fnOutput     = os.path.join(pp_dir, 'PPToolboxOutput.pptx')
fig_dir         = r'C:\Users\N.Srisutthiyakorn\OneDrive - Shell\Documents\GitHub\Python_Toolbox_NS\PPToolboxExample\MainFigureFolder'

createPPT(pp_fnTemplate, pp_fnOutput, fig_dir)   